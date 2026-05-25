"""
Auto-ingest script for Fristine Infotech document chatbot.
Runs via GitHub Actions when files are pushed to the repo.
Supports: PDF, PPTX, DOCX, TXT files.
New files are embedded and stored in Supabase. Already-indexed files are skipped.
"""

import os, json, time, urllib.request, urllib.parse

GEMINI_KEY = os.environ['GEMINI_API_KEY']
SUPA_URL   = os.environ['SUPABASE_URL'].rstrip('/')
SUPA_KEY   = os.environ['SUPABASE_KEY']

SUPA_HEADERS = {
    'Content-Type': 'application/json',
    'Authorization': f'Bearer {SUPA_KEY}',
    'apikey': SUPA_KEY,
}

# All files are attempted — binary/unreadable ones are skipped automatically.
# Known skip extensions (images, executables, archives, etc.)
SKIP_EXTS = {
    '.png', '.jpg', '.jpeg', '.gif', '.bmp', '.ico', '.svg', '.webp',
    '.mp3', '.mp4', '.wav', '.avi', '.mov', '.mkv',
    '.zip', '.tar', '.gz', '.rar', '.7z',
    '.exe', '.dll', '.so', '.bin', '.pyc',
    '.lock', '.sum', '.mod',
}

# ── Text extraction ────────────────────────────────────────────────────────────

def extract_pdf(path):
    import fitz
    doc = fitz.open(path)
    return '\n'.join(page.get_text() for page in doc)

def extract_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                parts.append(shape.text.strip())
    return '\n'.join(parts)

def extract_docx(path):
    from docx import Document
    doc = Document(path)
    return '\n'.join(p.text for p in doc.paragraphs if p.text.strip())

def extract_xlsx(path):
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True)
    rows = []
    for sheet in wb.worksheets:
        rows.append(f'[Sheet: {sheet.title}]')
        row_count = 0
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                rows.append(' | '.join(cells))
                row_count += 1
            if row_count >= 2000:  # cap at 2000 rows per sheet
                rows.append('[...truncated]')
                break
    return '\n'.join(rows)

def extract_csv(path):
    import csv
    rows = []
    with open(path, encoding='utf-8', errors='ignore', newline='') as f:
        reader = csv.reader(f)
        for row in reader:
            line = ' | '.join(cell.strip() for cell in row if cell.strip())
            if line:
                rows.append(line)
    return '\n'.join(rows)

def extract_html(path):
    from html.parser import HTMLParser
    class TextExtractor(HTMLParser):
        def __init__(self):
            super().__init__()
            self.parts = []
            self._skip = False
        def handle_starttag(self, tag, attrs):
            if tag in ('script', 'style'):
                self._skip = True
        def handle_endtag(self, tag):
            if tag in ('script', 'style'):
                self._skip = False
        def handle_data(self, data):
            if not self._skip and data.strip():
                self.parts.append(data.strip())
    with open(path, encoding='utf-8', errors='ignore') as f:
        raw = f.read()
    p = TextExtractor()
    p.feed(raw)
    return '\n'.join(p.parts)

def read_as_text(path):
    """Fallback: try reading the file as UTF-8 text."""
    try:
        with open(path, encoding='utf-8', errors='ignore') as f:
            return f.read()
    except Exception:
        return ''

def extract_text(path):
    ext = os.path.splitext(path)[1].lower()
    if ext in SKIP_EXTS:
        return ''
    if ext == '.pdf':
        return extract_pdf(path)
    elif ext in ('.pptx', '.ppt'):
        return extract_pptx(path)
    elif ext in ('.docx', '.doc'):
        return extract_docx(path)
    elif ext in ('.xlsx', '.xls'):
        return extract_xlsx(path)
    elif ext == '.csv':
        return extract_csv(path)
    elif ext in ('.html', '.htm'):
        return extract_html(path)
    else:
        # TXT, MD, JSON, RTF, YAML, XML, code files — all read as text
        return read_as_text(path)

# ── Chunking ───────────────────────────────────────────────────────────────────

def chunk_text(text, size=600, overlap=100):
    cleaned = text.replace('\r\n', '\n').strip()
    chunks = []
    start = 0
    while start < len(cleaned):
        end = min(start + size, len(cleaned))
        if end < len(cleaned):
            dot = cleaned.rfind('.', start + int(size * 0.4), end)
            nl  = cleaned.rfind('\n', start + int(size * 0.4), end)
            best = max(dot, nl)
            if best > start:
                end = best + 1
        chunk = cleaned[start:end].strip()
        if len(chunk) > 40:
            chunks.append(chunk)
        next_start = end - overlap
        if next_start <= start:
            next_start = start + max(1, size - overlap)
        start = next_start
    return chunks

# ── HTTP helpers ───────────────────────────────────────────────────────────────

def http_post(url, headers, data):
    body = json.dumps(data).encode()
    req = urllib.request.Request(url, data=body, headers=headers, method='POST')
    try:
        with urllib.request.urlopen(req) as r:
            return r.status, r.read().decode()
    except urllib.error.HTTPError as e:
        return e.code, e.read().decode()

def http_delete(url, headers):
    req = urllib.request.Request(url, headers=headers, method='DELETE')
    try:
        with urllib.request.urlopen(req) as r:
            r.read()
    except:
        pass

def get_indexed_sources():
    req = urllib.request.Request(
        f'{SUPA_URL}/rest/v1/documents?select=source',
        headers=SUPA_HEADERS
    )
    with urllib.request.urlopen(req) as r:
        rows = json.load(r)
    sources = set()
    for row in rows:
        sources.add(row['source'])
    return sources

def embed_text(text):
    url = f'https://generativelanguage.googleapis.com/v1beta/models/gemini-embedding-001:embedContent?key={GEMINI_KEY}'
    status, body = http_post(url, {'Content-Type': 'application/json'},
                              {'content': {'parts': [{'text': text}]}, 'outputDimensionality': 768})
    if status != 200:
        msg = json.loads(body).get('error', {}).get('message', f'HTTP {status}')
        raise Exception(f'[{status}] {msg}')
    return json.loads(body)['embedding']['values']

def insert_row(content, embedding, source, sha=None):
    row = {'content': content, 'embedding': f'[{",".join(str(v) for v in embedding)}]', 'source': source}
    if sha:
        row['sha'] = sha
    status, body = http_post(
        f'{SUPA_URL}/rest/v1/documents',
        {**SUPA_HEADERS, 'Prefer': 'return=minimal'},
        row
    )
    if status >= 300:
        raise Exception(f'Supabase {status}: {body}')

def delete_source(source):
    http_delete(
        f'{SUPA_URL}/rest/v1/documents?source=eq.{urllib.parse.quote(source)}',
        SUPA_HEADERS
    )

# ── Main ───────────────────────────────────────────────────────────────────────

def ingest_file(filepath, source_name, sha=None, force=False):
    print(f'\n📄 Processing: {source_name}')
    text = extract_text(filepath)
    if not text or len(text) < 50:
        print('  ⚠️  No text extracted, skipping.')
        return 0

    print(f'  Text length: {len(text)} chars')
    chunks = chunk_text(text)
    print(f'  Chunks: {len(chunks)}')

    if force:
        delete_source(source_name)

    ok = fail = 0
    for i, chunk in enumerate(chunks):
        try:
            emb = embed_text(chunk)
            insert_row(chunk, emb, source_name, sha=sha)
            ok += 1
            print(f'  [{i+1}/{len(chunks)}] ✓', flush=True)
            time.sleep(0.4)
        except Exception as e:
            print(f'  [{i+1}/{len(chunks)}] ✗ {e}')
            fail += 1
            if '403' in str(e) or 'leaked' in str(e).lower() or 'invalid' in str(e).lower():
                print('  API key error — stopping.')
                break
            elif '429' in str(e):
                print('  Rate limited — waiting 15s...')
                time.sleep(15)
            else:
                time.sleep(1)

    print(f'  Done: {ok} inserted, {fail} failed')
    return ok


def get_git_sha(filepath):
    """Get the git blob SHA for a file to detect changes."""
    import subprocess
    try:
        result = subprocess.run(['git', 'hash-object', filepath], capture_output=True, text=True)
        return result.stdout.strip()
    except:
        return None

def get_indexed_shas():
    """Get {source_name: sha} from Supabase metadata column."""
    req = urllib.request.Request(
        f'{SUPA_URL}/rest/v1/documents?select=source,sha&limit=1000',
        headers=SUPA_HEADERS
    )
    try:
        with urllib.request.urlopen(req) as r:
            rows = json.load(r)
        result = {}
        for row in rows:
            src = row.get('source')
            sha = row.get('sha')
            if src and src not in result:
                result[src] = sha
        return result
    except:
        return {}

def main():
    print('🔍 Scanning repo for documents...')
    indexed_shas = get_indexed_shas()
    print(f'Already indexed: {list(indexed_shas.keys()) if indexed_shas else "none"}')

    # Find all supported files in the repo
    all_files = []
    for root, dirs, files in os.walk('.'):
        dirs[:] = [d for d in dirs if d not in ('.github', 'scripts', '.git', 'node_modules', '__pycache__')]
        for fname in files:
                ext = os.path.splitext(fname)[1].lower()
                if ext not in SKIP_EXTS and not fname.startswith('.'):
                    all_files.append(os.path.join(root, fname))

    if not all_files:
        print('No supported files found in repo.')
        return

    print(f'\nFound {len(all_files)} file(s): {[os.path.basename(f) for f in all_files]}')

    total = 0
    for filepath in all_files:
        source_name = os.path.basename(filepath)
        current_sha = get_git_sha(filepath)
        stored_sha  = indexed_shas.get(source_name)

        if stored_sha and stored_sha == current_sha:
            print(f'\n⏭️  Unchanged, skipping: {source_name}')
            continue

        if stored_sha:
            print(f'\n🔄 File changed, re-indexing: {source_name}')
            action = 'updated'
        else:
            print(f'\n🆕 New file: {source_name}')
            action = 'added'

        inserted = ingest_file(filepath, source_name, sha=current_sha, force=True)
        if inserted > 0:
            log_policy_update(source_name, action)
        total += inserted

    print(f'\n✅ Ingest complete. Total chunks inserted/updated: {total}')


def log_policy_update(source, action):
    """Log to Supabase so the chatbot can show update alerts to employees."""
    try:
        http_post(
            f'{SUPA_URL}/rest/v1/policy_updates',
            {**SUPA_HEADERS, 'Prefer': 'return=minimal'},
            {'source': source, 'action': action}
        )
    except:
        pass


if __name__ == '__main__':
    main()
